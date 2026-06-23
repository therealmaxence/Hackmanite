import structlog

logger = structlog.get_logger()

def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            paragraphs.extend(cell.text for row in table.rows for cell in row.cells if cell.text.strip())
        return "\n".join(paragraphs)
    except Exception as e:
        logger.error("DOCX extraction failed", error=str(e), path=path)
        return ""

def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        return "\n".join(" ".join(str(c) for c in row if c is not None) for sheet in wb.worksheets for row in sheet.iter_rows(values_only=True) if " ".join(str(c) for c in row if c is not None).strip())
    except Exception as e:
        logger.error("XLSX extraction failed", error=str(e), path=path)
        return ""

def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        return "\n".join(shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
    except Exception as e:
        logger.error("PPTX extraction failed", error=str(e), path=path)
        return ""
